"""Tests for file format parsers."""
import tempfile
from pathlib import Path
from app.services.parsers.text_parser import TextParser


class TestTextParser:
    def test_supports_txt(self):
        assert TextParser.supports("text/plain") is True

    def test_parse_utf8(self):
        content = "这是广告文案内容\n包含多行文本"
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".txt", delete=False, encoding="utf-8"
        ) as f:
            f.write(content)
            tmp_path = Path(f.name)

        try:
            result = TextParser.parse(tmp_path)
            # Normalize line endings for cross-platform compatibility
            assert result.text.replace("\r\n", "\n").strip() == content.strip()
            assert result.source_format == "txt_raw"
            assert result.quality == "good"
            assert result.fallback_used is False
        finally:
            tmp_path.unlink()

    def test_parse_gbk(self):
        content = "GBK编码的广告文案"
        with tempfile.NamedTemporaryFile(
            mode="wb", suffix=".txt", delete=False
        ) as f:
            f.write(content.encode("gbk"))
            tmp_path = Path(f.name)

        try:
            result = TextParser.parse(tmp_path)
            assert content in result.text
            assert result.quality == "good"
        finally:
            tmp_path.unlink()

    def test_parse_empty_file(self):
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".txt", delete=False, encoding="utf-8"
        ) as f:
            tmp_path = Path(f.name)

        try:
            result = TextParser.parse(tmp_path)
            assert result.text == ""
            assert result.quality == "minimal"
        finally:
            tmp_path.unlink()


class TestImageParser:
    def test_supports_image_mimes(self):
        from app.services.parsers.image_parser import ImageParser
        assert ImageParser.supports("image/png") is True
        assert ImageParser.supports("image/jpeg") is True
        assert ImageParser.supports("image/gif") is True
        assert ImageParser.supports("image/bmp") is True
        assert ImageParser.supports("video/mp4") is False

    def test_parse_creates_result(self):
        from app.services.parsers.image_parser import ImageParser
        from PIL import Image
        import tempfile
        from pathlib import Path

        img = Image.new("RGB", (200, 50), color="white")
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        img.save(tmp)
        tmp.close()
        tmp_path = Path(tmp.name)

        try:
            result = ImageParser.parse(tmp_path)
            assert result.source_format == "image_ocr"
            assert result.quality in ("good", "degraded", "minimal")
            assert isinstance(result.text, str)
        finally:
            try:
                tmp_path.unlink()
            except FileNotFoundError:
                pass


class TestOfficeParser:
    def test_supports_docx(self):
        from app.services.parsers.office_parser import OfficeParser
        assert OfficeParser.supports("application/vnd.openxmlformats-officedocument.wordprocessingml.document") is True

    def test_supports_pptx(self):
        from app.services.parsers.office_parser import OfficeParser
        assert OfficeParser.supports("application/vnd.openxmlformats-officedocument.presentationml.presentation") is True

    def test_supports_xlsx(self):
        from app.services.parsers.office_parser import OfficeParser
        assert OfficeParser.supports("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet") is True

    def test_supports_doc(self):
        from app.services.parsers.office_parser import OfficeParser
        # Legacy binary .doc files cannot be parsed safely by python-docx.
        assert OfficeParser.supports("application/msword") is False

    def test_parse_docx(self):
        from app.services.parsers.office_parser import OfficeParser
        from docx import Document
        import tempfile
        from pathlib import Path

        tmp = tempfile.NamedTemporaryFile(suffix=".docx", delete=False)
        tmp.close()
        doc = Document()
        doc.add_paragraph("广告文案第一段")
        doc.add_paragraph("广告文案第二段")
        doc.save(tmp.name)
        tmp_path = Path(tmp.name)

        try:
            result = OfficeParser.parse(tmp_path, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
            assert "广告文案第一段" in result.text
            assert "广告文案第二段" in result.text
            assert result.source_format == "docx_parse"
            assert result.quality == "good"
        finally:
            tmp_path.unlink()

    def test_parse_pptx(self):
        from app.services.parsers.office_parser import OfficeParser
        from pptx import Presentation
        import tempfile
        from pathlib import Path

        tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        tmp.close()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "广告标题"
        prs.save(tmp.name)
        tmp_path = Path(tmp.name)

        try:
            result = OfficeParser.parse(tmp_path, "application/vnd.openxmlformats-officedocument.presentationml.presentation")
            assert "广告标题" in result.text
            assert result.source_format == "pptx_parse"
        finally:
            tmp_path.unlink()

    def test_parse_xlsx(self):
        from app.services.parsers.office_parser import OfficeParser
        from openpyxl import Workbook
        import tempfile
        from pathlib import Path

        tmp = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False)
        tmp.close()
        wb = Workbook()
        ws = wb.active
        ws["A1"] = "产品名称"
        ws["B1"] = "广告文案"
        ws["A2"] = "产品A"
        ws["B2"] = "优质产品宣传语"
        wb.save(tmp.name)
        tmp_path = Path(tmp.name)

        try:
            result = OfficeParser.parse(tmp_path, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
            assert "产品名称" in result.text
            assert "优质产品宣传语" in result.text
            assert result.source_format == "xlsx_parse"
        finally:
            tmp_path.unlink()


class TestPdfParser:
    def test_supports_pdf(self):
        from app.services.parsers.pdf_parser import PdfParser
        assert PdfParser.supports("application/pdf") is True

    def test_parse_text_pdf(self):
        from app.services.parsers.pdf_parser import PdfParser
        import fitz
        import tempfile
        from pathlib import Path

        tmp = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
        tmp.close()
        doc = fitz.open()
        page = doc.new_page()
        text = "Advertising compliance review text with enough words to meet the quality threshold."
        page.insert_text((50, 50), text, fontsize=12)
        doc.save(tmp.name)
        doc.close()
        tmp_path = Path(tmp.name)

        try:
            result = PdfParser.parse(tmp_path)
            assert "Advertising compliance review text" in result.text
            assert result.source_format == "pdf_text"
            assert result.quality == "good"
            assert result.fallback_used is False
        finally:
            tmp_path.unlink()

    def test_parse_empty_pdf_falls_back_to_ocr(self):
        from app.services.parsers.pdf_parser import PdfParser
        import fitz
        import tempfile
        from pathlib import Path

        tmp = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
        tmp.close()
        doc = fitz.open()
        doc.new_page()
        doc.save(tmp.name)
        doc.close()
        tmp_path = Path(tmp.name)

        try:
            result = PdfParser.parse(tmp_path)
            # OCR fallback returns pdf_ocr regardless of whether
            # tesseract is installed (empty text if not available)
            assert result.source_format == "pdf_ocr"
            assert result.fallback_used is True
        finally:
            tmp_path.unlink()
