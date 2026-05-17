"""Office document parser — DOCX, DOC, PPTX, XLSX text extraction."""
from __future__ import annotations

from pathlib import Path

from app.services.parsers.text_parser import ExtractionResult


class OfficeParser:
    MIME_TYPES = {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    }

    @staticmethod
    def supports(mime: str) -> bool:
        return mime in OfficeParser.MIME_TYPES

    @staticmethod
    def parse(file_path: str | Path, mime: str) -> ExtractionResult:
        path = Path(file_path)
        if "wordprocessingml" in mime:
            return _parse_docx(path)
        elif "presentation" in mime:
            return _parse_pptx(path)
        elif "spreadsheet" in mime:
            return _parse_xlsx(path)
        elif "msword" in mime:
            return _parse_docx(path)
        raise ValueError(f"Unsupported office MIME: {mime}")


def _parse_docx(path: Path) -> ExtractionResult:
    from docx import Document

    doc = Document(str(path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    text = "\n".join(paragraphs)
    quality = "good" if text else "minimal"
    return ExtractionResult(
        text=text,
        source_format="docx_parse",
        quality=quality,
        fallback_used=False,
    )


def _parse_pptx(path: Path) -> ExtractionResult:
    from pptx import Presentation

    prs = Presentation(str(path))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
    text = "\n".join(texts)
    quality = "good" if text else "minimal"
    return ExtractionResult(
        text=text,
        source_format="pptx_parse",
        quality=quality,
        fallback_used=False,
    )


def _parse_xlsx(path: Path) -> ExtractionResult:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    texts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                texts.append(row_text)
    wb.close()
    text = "\n".join(texts)
    quality = "good" if text else "minimal"
    return ExtractionResult(
        text=text,
        source_format="xlsx_parse",
        quality=quality,
        fallback_used=False,
    )
